import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import subprocess
import os
import streamlit as st
import ollama
import tempfile
import random

# Run this script with: streamlit run ppt_generator.py
# Do not run directly with python ppt_generator.py to avoid 'missing ScriptRunContext' warnings.

# Function to interact with LLaMA with CSV-specific content
def generate_with_llama(prompt):
    try:
        response = ollama.generate(model="llama3.2", prompt=f"{prompt} Provide only the concise, complete text or numbered list (no introductory phrases, no formatting). Ensure 5 to 6 complete bullet points ending with full sentences, derived solely from the provided CSV data analysis.")
        return response['response'].strip()
    except Exception:
        return "Analysis failed due to error.\nCSV data could not be processed.\nPlease verify file integrity.\nContact support for assistance.\nThis is an error state."

# Function to split text into 5 to 6 bullet points
def split_into_bullets(text, min_points=5, max_points=6):
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    if not lines or len(lines) < min_points:
        return [
            "Insufficient data in CSV.",
            "Analysis cannot be completed.",
            "Please check CSV content.",
            "No insights can be derived.",
            "This is an error message."
        ]
    num_points = random.randint(min_points, min(max_points, len(lines)))
    return lines[:num_points]

# Function to add a slide with adjusted layout
def add_slide(prs, title, content=None, chart_path=None, bg_color=RGBColor(240, 240, 240)):
    slide_layout = prs.slide_layouts[5] if chart_path else prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_shape.top = Inches(0.5)
    title_shape.left = Inches(1)
    title_shape.width = Inches(8)
    title_shape.height = Inches(1)
    
    # Content or plot
    if chart_path:
        slide.shapes.add_picture(chart_path, Inches(1), Inches(1.75), Inches(8), Inches(5))
    elif content:
        textbox_height = Inches(6)
        font_size = Pt(14)
        tf = slide.placeholders[1].text_frame if slide.placeholders else slide.shapes.add_textbox(Inches(1), Inches(1.75), Inches(8), textbox_height).text_frame
        tf.clear()
        for point in content:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = font_size
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.level = 0
            p.space_after = Pt(6)
    return slide

# Backend EDA and slide generation
def generate_eda_report(csv_file, col, plot_type, min_slides, user_prompt):
    csv_file.seek(0)
    try:
        df = pd.read_csv(csv_file)
        if df.empty:
            return False, "CSV file is empty."
    except Exception as e:
        return False, f"Error reading CSV: {str(e)}"
    
    num_cols = len(df.columns)
    other_cols = [c for c in df.columns if c != col]
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # First slide: Only title
    title_prompt = f"Analyze CSV: Rows={len(df)}, Cols={num_cols}, Selected={col}. Generate a 5-word title based on data and '{user_prompt}'."
    cover_title = generate_with_llama(title_prompt).split('\n')[0]
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(240, 248, 255)
    title = slide.shapes.title
    title.text = cover_title
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title.top = Inches(3)
    title.left = Inches(1)
    title.width = Inches(8)
    for shape in slide.shapes:
        if shape.placeholder_format.idx != 0:
            shape.element.getparent().remove(shape.element)
    slide_titles = [cover_title]
    
    # Second slide(s): Overview of upcoming slides
    slide_titles.extend(["Introduction to Analysis"])
    slide_titles.extend([f"Comparison Plot: {col} vs {other_col}" for other_col in other_cols])
    slide_titles.extend([f"Comparison Insights: {col} vs {other_col}" for other_col in other_cols])
    slide_titles.extend([f"Detailed Insights: {col} vs {other_col}" for other_col in other_cols])
    extra_slides_needed = min_slides > (2 * num_cols)  # Define here for later use
    if extra_slides_needed:
        slide_titles.extend(["Index of Slides"])
    if user_prompt.lower() != "default analysis of one column vs others" and "summary" in user_prompt.lower():
        slide_titles.append("Summary of Findings")
    slide_titles.append("Conclusion of Analysis")
    overview_content = [f"{i + 1}. {title}" for i, title in enumerate(slide_titles[2:-1])]  # Skip title, overview, and Thank You
    max_points_per_slide = 6  # Limit to 6 to fit
    for i in range(0, len(overview_content), max_points_per_slide):
        chunk = overview_content[i:i + max_points_per_slide]
        title = "Overview of Upcoming Slides" if i == 0 else "Overview of Upcoming Slides Continued"
        add_slide(prs, title, chunk)
    
    # Add introduction slide (always after overview)
    intro_prompt = f"Introduce analysis of {col} vs others based on CSV with {len(df)} rows, {num_cols} columns, focusing on {col} in 5 to 6 bullet points based on '{user_prompt}'."
    intro_text = generate_with_llama(intro_prompt)
    intro_points = split_into_bullets(intro_text)
    add_slide(prs, "Introduction to Analysis", intro_points)
    
    # Generate comparison slides with CSV-specific content
    for idx, other_col in enumerate(other_cols):
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            chart_path = tmp.name
            plt.figure(figsize=(8, 5))
            is_numeric_col = pd.api.types.is_numeric_dtype(df[col])
            is_numeric_other = pd.api.types.is_numeric_dtype(df[other_col])
            actual_plot_type = plot_type
            
            if is_numeric_col and is_numeric_other:
                if plot_type == "Scatter":
                    df.plot.scatter(x=col, y=other_col, color='teal', alpha=0.5)
                    plt.title(f"{col} vs {other_col}", fontsize=12)
                elif plot_type == "Hexbin":
                    plt.hexbin(df[col], df[other_col], gridsize=20, cmap='Blues', mincnt=1)
                    plt.colorbar(label='Count')
                    plt.title(f"{col} vs {other_col}", fontsize=12)
                elif plot_type == "Box":
                    bins = pd.cut(df[col], bins=3)
                    df_box = df[[other_col]].copy()
                    df_box['Binned_' + col] = bins
                    df_box.boxplot(column=other_col, by='Binned_' + col, grid=False, patch_artist=True)
                    plt.title(f"{other_col} by {col}", fontsize=12)
                    plt.suptitle('')
                elif plot_type == "Bar":
                    bins = pd.cut(df[col], bins=3)
                    df.groupby(bins)[other_col].mean().plot(kind='bar', color='lightcoral')
                    plt.title(f"Mean {other_col} by {col}", fontsize=12)
            elif is_numeric_other:
                actual_plot_type = "Bar"
                df.groupby(col)[other_col].mean().plot(kind='bar', color='lightgreen')
                plt.title(f"Mean {other_col} by {col}", fontsize=12)
            elif is_numeric_col:
                actual_plot_type = "Bar"
                df.groupby(other_col)[col].count().plot(kind='bar', color='lightblue')
                plt.title(f"Count of {col} by {other_col}", fontsize=12)
            else:
                actual_plot_type = "Stacked Bar"
                pd.crosstab(df[col], df[other_col]).plot(kind='bar', stacked=True, colormap='Set2')
                plt.title(f"{col} vs {other_col}", fontsize=12)
            
            plt.xlabel(col, fontsize=10)
            plt.ylabel(other_col, fontsize=10)
            plt.xticks(rotation=45, ha='right', fontsize=8)
            plt.savefig(chart_path, bbox_inches='tight')
            plt.close()
        
        # Plot slide (only plot)
        add_slide(prs, f"Comparison Plot: {col} vs {other_col}", chart_path=chart_path)
        
        # Content slide (CSV-specific stats)
        corr = df[[col, other_col]].corr().iloc[0,1] if is_numeric_col and is_numeric_other else "N/A"
        col_stat_label = "mean" if is_numeric_col else "unique"
        col_stat_value = f"{df[col].mean():.2f}" if is_numeric_col else str(df[col].nunique())
        other_col_stat_label = "mean" if is_numeric_other else "unique"
        other_col_stat_value = f"{df[other_col].mean():.2f}" if is_numeric_other else str(df[other_col].nunique())
        stats_content = f"{col} vs {other_col}: Corr={corr if corr != 'N/A' else 'N/A'}, {col} {col_stat_label}={col_stat_value}, {other_col} {other_col_stat_label}={other_col_stat_value}"
        content_points = [
            f"Rows analyzed: {len(df)}. Total entries in CSV.",
            f"Correlation: {corr:.2f}. Shows {col} vs {other_col} link." if corr != "N/A" else f"{col} type: {df[col].dtype}. Non-numeric data detected.",
            f"{other_col} mean: {df[other_col].mean():.2f}. Average from CSV data." if is_numeric_other else f"{other_col} unique: {df[other_col].nunique()}. Distinct values counted.",
            f"{other_col} min: {df[other_col].min():.2f}. Minimum value in CSV." if is_numeric_other else f"{other_col} top: {df[other_col].mode()[0]}. Most frequent in CSV.",
            f"{other_col} max: {df[other_col].max():.2f}. Maximum value in CSV." if is_numeric_other else f"{other_col} diversity: {'High' if df[other_col].nunique() > 5 else 'Low'}. Variation in data.",
            f"{col} stat: {df[col].mean():.2f}. Numeric average from CSV." if is_numeric_col else f"{col} top: {df[col].mode()[0]}. Top category in CSV."
        ]
        add_slide(prs, f"Comparison Insights: {col} vs {other_col}", content_points)
        
        # Detailed insights slide (CSV-specific)
        detail_prompt = f"Provide detailed insights for {col} vs {other_col} based on CSV data: '{stats_content}', in 5 to 6 bullet points based on '{user_prompt}'."
        detail_text = generate_with_llama(detail_prompt)
        detail_points = split_into_bullets(detail_text)
        add_slide(prs, f"Detailed Insights: {col} vs {other_col}", detail_points)
        
        os.remove(chart_path)
    
    # Add index slide if needed
    if extra_slides_needed:
        index_content = [f"{i + 1}. {title}" for i, title in enumerate(slide_titles[2:-1])]  # Skip title, overview, and Thank You
        index_points = split_into_bullets("\n".join(index_content))
        add_slide(prs, "Index of Slides", index_points)
    
    # Handle user prompt customization
    if user_prompt and user_prompt.lower() != "default analysis of one column vs others" and "summary" in user_prompt.lower():
        summary_prompt = f"Summarize analysis of {col} vs others based on CSV data with {len(df)} rows, {num_cols} columns in 5 to 6 bullet points based on '{user_prompt}'."
        summary_text = generate_with_llama(summary_prompt)
        summary_points = split_into_bullets(summary_text)
        add_slide(prs, "Summary of Findings", summary_points)
        slide_titles.append("Summary of Findings")
    
    # Ensure minimum slides by adding extra content slides if needed
    current_slides = len(slide_titles) + 1  # +1 for Thank You
    if current_slides < min_slides:
        for i in range(min_slides - current_slides):
            extra_prompt = f"Provide extra analysis for {col} vs others based on CSV data with {len(df)} rows, {num_cols} columns in 5 to 6 bullet points based on '{user_prompt}'."
            extra_text = generate_with_llama(extra_prompt)
            extra_points = split_into_bullets(extra_text)
            add_slide(prs, f"Additional Analysis {i + 1}", extra_points)
            slide_titles.append(f"Additional Analysis {i + 1}")
    
    # Second-to-last slide: Conclusion (CSV-specific)
    conclusion_prompt = f"Conclude analysis of {col} vs others based on CSV data with {len(df)} rows, {num_cols} columns in 5 to 6 bullet points based on '{user_prompt}'."
    conclusion_text = generate_with_llama(conclusion_prompt)
    conclusion_points = split_into_bullets(conclusion_text)
    add_slide(prs, "Conclusion of Analysis", conclusion_points)
    
    # Last slide: Thank You (only text, centered)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(240, 248, 255)
    title = slide.shapes.title
    title.text = "Thank You"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title.top = Inches(3)
    title.left = Inches(1)
    title.width = Inches(8)
    for shape in slide.shapes:
        if shape.placeholder_format.idx != 0:
            shape.element.getparent().remove(shape.element)
    
    # Save and convert
    pptx_file = "one_column_eda_report.pptx"
    odp_file = "one_column_eda_report.odp"
    prs.save(pptx_file)
    try:
        subprocess.run(["libreoffice", "--headless", "--convert-to", "odp", pptx_file], check=True)
        with open(odp_file, "rb") as f:
            odp_bytes = f.read()
        os.remove(pptx_file)
        os.remove(odp_file)
        return True, odp_bytes
    except Exception as e:
        return False, f"Error converting to ODP: {str(e)}"

# Streamlit front-end
def main():
    st.title("AI Based PPT Generator")
    st.markdown("Upload a CSV, select one column, choose a plot type, set minimum slides, and optionally provide a prompt.")
    
    uploaded_file = st.file_uploader("Choose a CSV file", type="csv")
    
    if uploaded_file:
        uploaded_file.seek(0)
        try:
            df = pd.read_csv(uploaded_file)
            if df.empty:
                st.error("CSV file is empty.")
                return
            columns = list(df.columns)
            st.success(f"Loaded {uploaded_file.name} with {len(df)} rows and {len(columns)} columns.")
        except Exception as e:
            st.error(f"Error loading CSV: {str(e)}")
            return
        
        col = st.selectbox("Select Column to Analyze", columns)
        plot_type = st.selectbox("Select Plot Type", ["Scatter", "Hexbin", "Box", "Bar"])
        min_slides = st.number_input("Minimum Number of Slides", min_value=3, value=5, step=1)  # Min 3 for title, overview, thank you
        user_prompt = st.text_area("Optional: Customize PPT (e.g., 'add summary slide')", 
                                   "Default analysis of one column vs others", height=100)
        
        if st.button("Generate Report"):
            with st.spinner("Generating report with LLaMA..."):
                success, result = generate_eda_report(uploaded_file, col, plot_type, min_slides, user_prompt)
                if success:
                    st.success("Report generated successfully!")
                    st.download_button(
                        label="Download EDA Report",
                        data=result,
                        file_name="one_column_eda_report.odp",
                        mime="application/vnd.oasis.opendocument.presentation"
                    )
                else:
                    st.error(f"Error: {result}")

if __name__ == "__main__":
    main()